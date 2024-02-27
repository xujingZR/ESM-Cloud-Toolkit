from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN

def add_2_pics(img_path_1, img_path_2, bullet_text, figcaption_1, figcaption_2, prs):

    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = bullet_text

    slide.shapes.add_picture(img_path_1, Cm(1.45), Cm(6.17), height=Cm(8.44))
    slide.shapes.add_picture(img_path_2, Cm(12.7), Cm(6.17), height=Cm(8.44))

    txBox = slide.shapes.add_textbox(Cm(2.65), Cm(14.89), Cm(8.86), Cm(2.18))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.add_paragraph()
    p.text = figcaption_1
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(17)

    txBox_2 = slide.shapes.add_textbox(Cm(13.9), Cm(14.89), Cm(8.86), Cm(2.18))
    tf_2 = txBox_2.text_frame
    tf_2.word_wrap = True

    p_2 = tf_2.add_paragraph()
    p_2.text = figcaption_2
    p_2.alignment = PP_ALIGN.CENTER
    p_2.font.size = Pt(17)

    return prs

def add_1_pic(img_path, bullet_text, figcaption, prs):

    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = bullet_text

    slide.shapes.add_picture(img_path, Cm(4.57), Cm(6.1), height=Cm(12.19))

    txBox = slide.shapes.add_textbox(Cm(2.65), Cm(14.89), Cm(8.86), Cm(2.18))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.add_paragraph()
    p.text = figcaption
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(17)

return prs

 
